from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import addphoto
import re
#Create a new PowerPoint presentation


def add_slide(prs, layout, title, subtitle):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title.strip()
    slide.placeholders[1].text=subtitle
    font = slide.shapes.title.text_frame.paragraphs[0].font
    font.name = 'Arial'
    font.size = Pt(30)
    font.bold = True
    font.italic = False
    for x in  slide.placeholders[1].text_frame.paragraphs:
        font1= x.font
        font1.name = 'Arial'
        font1.size = Pt(16)
        font1.bold = False
        font1.italic = False
    return slide

def add_slide_img(prs, layout, img_path):
    slide = prs.slides.add_slide(layout)
    img_path =  ""+img_path
    left =  Inches(1.10)
    top = Inches(0.7)
    width = Inches(8)
    height = Inches(6)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)
        
def presentate(gpt_response_object):
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[1]
    title_slide_laying=prs.slide_layouts[6]

    topic = gpt_response_object['Topic']
    summary = gpt_response_object['Summary']

    summary_length = len(summary)
    first_half_item = ' '.join(summary[0:int(summary_length/2)])
    second_half_item = ' '.join(summary[int(summary_length/2):])

    add_slide(prs, title_slide_layout,topic,first_half_item)
    add_slide(prs, title_slide_layout, topic,second_half_item)

    images = addphoto.get_images(topic, 2)
    if images:
        try:
            add_slide_img(prs, title_slide_laying, f"images/{images[0]}")
        except Exception as e:
            print(f"Failed to add first image: {e}")
    if len(images) > 1:
        try:
            add_slide_img(prs, title_slide_laying, f"images/{images[1]}")
        except Exception as e:
            print(f"Failed to add second image: {e}")
    # Save the presentation
    prs.save("./PPT.pptx")